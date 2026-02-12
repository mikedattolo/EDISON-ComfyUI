"""
Professional File Generation Module
Generates high-quality PDFs, HTML documents, slideshows, spreadsheets, and more
without external dependencies beyond Python stdlib.
"""

import json
import re
import os
import io
import textwrap
import html
import time
import struct
import zlib
from datetime import datetime
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
import logging

logger = logging.getLogger(__name__)


# ──────────────────────────────────────────────────────────────────
#  PDF Generation — Professional multi-page with headers, footers,
#  styled sections, bullet points, and proper typography
# ──────────────────────────────────────────────────────────────────

class PDFDocument:
    """Generate professional PDF documents from structured content."""

    FONTS = {
        "Helvetica": "Helvetica",
        "Helvetica-Bold": "Helvetica-Bold",
        "Helvetica-Oblique": "Helvetica-Oblique",
        "Courier": "Courier",
    }
    PAGE_W = 612   # Letter width in points
    PAGE_H = 792   # Letter height in points
    MARGIN_L = 72  # 1 inch
    MARGIN_R = 72
    MARGIN_T = 72
    MARGIN_B = 72
    LINE_H = 16    # Normal line height
    TITLE_H = 28   # Title line height
    H1_H = 22
    H2_H = 18

    def __init__(self, title: str = "Document", author: str = "EDISON AI"):
        self.title = title
        self.author = author
        self.objects: list = []
        self.pages: list = []  # Each page is a list of content stream lines
        self.current_page: list = []
        self.y = self.PAGE_H - self.MARGIN_T
        self.page_num = 0
        self.total_pages = 0
        self._font_objs: Dict[str, int] = {}

    def _new_page(self):
        if self.current_page:
            self.pages.append(self.current_page)
        self.current_page = []
        self.y = self.PAGE_H - self.MARGIN_T
        self.page_num += 1

    def _ensure_space(self, needed: float):
        if self.y - needed < self.MARGIN_B:
            self._new_page()

    @staticmethod
    def _escape(text: str) -> str:
        return text.replace("\\", "\\\\").replace("(", "\\(").replace(")", "\\)")

    def _text_line(self, text: str, x: float, y: float, font: str = "Helvetica",
                   size: float = 11, color: Tuple[float, float, float] = (0, 0, 0)):
        r, g, b = color
        safe = self._escape(text)
        self.current_page.append(f"{r:.2f} {g:.2f} {b:.2f} rg")
        self.current_page.append(f"BT /{font} {size} Tf {x:.1f} {y:.1f} Td ({safe}) Tj ET")

    def _draw_line(self, x1: float, y1: float, x2: float, y2: float,
                   width: float = 0.5, color: Tuple[float, float, float] = (0.7, 0.7, 0.7)):
        r, g, b = color
        self.current_page.append(f"{r:.2f} {g:.2f} {b:.2f} RG")
        self.current_page.append(f"{width:.1f} w")
        self.current_page.append(f"{x1:.1f} {y1:.1f} m {x2:.1f} {y2:.1f} l S")

    def _draw_rect(self, x: float, y: float, w: float, h: float,
                   fill: Tuple[float, float, float] = (0.95, 0.95, 0.97)):
        r, g, b = fill
        self.current_page.append(f"{r:.2f} {g:.2f} {b:.2f} rg")
        self.current_page.append(f"{x:.1f} {y:.1f} {w:.1f} {h:.1f} re f")

    def add_title(self, text: str):
        self._ensure_space(self.TITLE_H + 20)
        # Title background accent bar
        self._draw_rect(self.MARGIN_L - 4, self.y - 6, 4, self.TITLE_H + 4, (0.26, 0.49, 0.92))
        self._text_line(text, self.MARGIN_L + 8, self.y, "Helvetica-Bold", 22, (0.1, 0.1, 0.15))
        self.y -= self.TITLE_H + 12
        # Underline
        self._draw_line(self.MARGIN_L, self.y + 6, self.PAGE_W - self.MARGIN_R, self.y + 6, 1.0, (0.26, 0.49, 0.92))
        self.y -= 8

    def add_subtitle(self, text: str):
        self._ensure_space(self.H2_H + 8)
        self._text_line(text, self.MARGIN_L + 8, self.y, "Helvetica-Oblique", 13, (0.4, 0.4, 0.5))
        self.y -= self.H2_H

    def add_heading(self, text: str, level: int = 1):
        h = self.H1_H if level == 1 else self.H2_H
        size = 16 if level == 1 else 13
        self._ensure_space(h + 16)
        self.y -= 8  # Extra space before heading
        if level == 1:
            self._draw_rect(self.MARGIN_L - 2, self.y - 3, 3, h - 2, (0.26, 0.49, 0.92))
        self._text_line(text, self.MARGIN_L + 6, self.y, "Helvetica-Bold", size, (0.15, 0.15, 0.2))
        self.y -= h

    def add_paragraph(self, text: str):
        usable_w = self.PAGE_W - self.MARGIN_L - self.MARGIN_R
        # Approximate chars per line at 11pt Helvetica (~6pt per char)
        chars_per_line = int(usable_w / 6.2)
        wrapped = textwrap.wrap(text, width=chars_per_line)
        for line in wrapped:
            self._ensure_space(self.LINE_H)
            self._text_line(line, self.MARGIN_L, self.y, "Helvetica", 11, (0.15, 0.15, 0.2))
            self.y -= self.LINE_H
        self.y -= 6  # Paragraph spacing

    def add_bullet(self, text: str, indent: int = 0):
        usable_w = self.PAGE_W - self.MARGIN_L - self.MARGIN_R - 20 - (indent * 15)
        chars_per_line = int(usable_w / 6.2)
        wrapped = textwrap.wrap(text, width=chars_per_line)
        x_base = self.MARGIN_L + 12 + (indent * 15)
        for i, line in enumerate(wrapped):
            self._ensure_space(self.LINE_H)
            if i == 0:
                bullet = "•" if indent == 0 else "◦"
                self._text_line(bullet, x_base - 10, self.y, "Helvetica", 11, (0.26, 0.49, 0.92))
            self._text_line(line, x_base, self.y, "Helvetica", 11, (0.15, 0.15, 0.2))
            self.y -= self.LINE_H
        self.y -= 2

    def add_numbered_item(self, number: int, text: str):
        usable_w = self.PAGE_W - self.MARGIN_L - self.MARGIN_R - 30
        chars_per_line = int(usable_w / 6.2)
        wrapped = textwrap.wrap(text, width=chars_per_line)
        for i, line in enumerate(wrapped):
            self._ensure_space(self.LINE_H)
            if i == 0:
                self._text_line(f"{number}.", self.MARGIN_L + 4, self.y, "Helvetica-Bold", 11, (0.26, 0.49, 0.92))
            self._text_line(line, self.MARGIN_L + 24, self.y, "Helvetica", 11, (0.15, 0.15, 0.2))
            self.y -= self.LINE_H
        self.y -= 2

    def add_code_block(self, code: str, language: str = ""):
        lines = code.splitlines()
        block_h = len(lines) * 13 + 16
        self._ensure_space(min(block_h, 200))  # At least try to fit some
        # Background
        self._draw_rect(self.MARGIN_L, self.y - block_h + 10, 
                       self.PAGE_W - self.MARGIN_L - self.MARGIN_R, block_h, 
                       (0.95, 0.95, 0.97))
        self.y -= 8
        for line in lines:
            self._ensure_space(13)
            safe_line = line[:100]  # Truncate long lines
            self._text_line(safe_line, self.MARGIN_L + 8, self.y, "Courier", 9, (0.2, 0.2, 0.25))
            self.y -= 13
        self.y -= 8

    def add_table(self, headers: List[str], rows: List[List[str]]):
        n_cols = len(headers)
        usable_w = self.PAGE_W - self.MARGIN_L - self.MARGIN_R
        col_w = usable_w / n_cols
        row_h = 20

        # Header background
        self._ensure_space(row_h * (len(rows) + 2))
        self._draw_rect(self.MARGIN_L, self.y - row_h + 4, usable_w, row_h, (0.22, 0.42, 0.82))
        for i, h in enumerate(headers):
            truncated = h[:int(col_w / 6)]
            self._text_line(truncated, self.MARGIN_L + i * col_w + 4, self.y - 10,
                          "Helvetica-Bold", 10, (1, 1, 1))
        self.y -= row_h

        # Data rows
        for r_idx, row in enumerate(rows):
            self._ensure_space(row_h)
            if r_idx % 2 == 0:
                self._draw_rect(self.MARGIN_L, self.y - row_h + 4, usable_w, row_h, (0.96, 0.97, 0.99))
            for i, cell in enumerate(row[:n_cols]):
                truncated = str(cell)[:int(col_w / 6)]
                self._text_line(truncated, self.MARGIN_L + i * col_w + 4, self.y - 10,
                              "Helvetica", 10, (0.2, 0.2, 0.25))
            self.y -= row_h
            # Row border
            self._draw_line(self.MARGIN_L, self.y + 4, self.MARGIN_L + usable_w, self.y + 4, 0.3)
        self.y -= 8

    def add_spacer(self, height: float = 12):
        self.y -= height

    def add_callout(self, text: str, style: str = "info"):
        """Add a highlighted callout/info box."""
        colors = {
            "info": (0.91, 0.95, 1.0),
            "warning": (1.0, 0.97, 0.89),
            "success": (0.91, 0.98, 0.93),
            "error": (1.0, 0.92, 0.92),
        }
        accent_colors = {
            "info": (0.26, 0.49, 0.92),
            "warning": (0.85, 0.65, 0.13),
            "success": (0.06, 0.73, 0.51),
            "error": (0.94, 0.27, 0.27),
        }
        fill = colors.get(style, colors["info"])
        accent = accent_colors.get(style, accent_colors["info"])

        usable_w = self.PAGE_W - self.MARGIN_L - self.MARGIN_R
        chars_per_line = int((usable_w - 20) / 6.2)
        wrapped = textwrap.wrap(text, width=chars_per_line)
        box_h = len(wrapped) * self.LINE_H + 16

        self._ensure_space(box_h + 8)
        # Background
        self._draw_rect(self.MARGIN_L, self.y - box_h + 6, usable_w, box_h, fill)
        # Accent bar
        self._draw_rect(self.MARGIN_L, self.y - box_h + 6, 4, box_h, accent)
        self.y -= 8
        for line in wrapped:
            self._text_line(line, self.MARGIN_L + 14, self.y, "Helvetica", 10, (0.2, 0.2, 0.25))
            self.y -= self.LINE_H
        self.y -= 8

    def _add_page_furniture(self, page_content: list, page_num: int, total_pages: int) -> list:
        """Add header/footer to a page's content stream."""
        furniture = []
        # Header line
        furniture.append(f"0.85 0.85 0.88 RG")
        furniture.append(f"0.5 w")
        furniture.append(f"{self.MARGIN_L} {self.PAGE_H - 52} m {self.PAGE_W - self.MARGIN_R} {self.PAGE_H - 52} l S")

        # Header text
        safe_title = self._escape(self.title[:60])
        furniture.append(f"0.55 0.55 0.6 rg")
        furniture.append(f"BT /Helvetica 8 Tf {self.MARGIN_L} {self.PAGE_H - 46} Td ({safe_title}) Tj ET")

        date_str = datetime.now().strftime("%B %d, %Y")
        safe_date = self._escape(date_str)
        furniture.append(f"BT /Helvetica 8 Tf {self.PAGE_W - self.MARGIN_R - 80} {self.PAGE_H - 46} Td ({safe_date}) Tj ET")

        # Footer line
        furniture.append(f"{self.MARGIN_L} {self.MARGIN_B - 10} m {self.PAGE_W - self.MARGIN_R} {self.MARGIN_B - 10} l S")

        # Page number
        page_text = f"Page {page_num} of {total_pages}"
        safe_page = self._escape(page_text)
        furniture.append(f"BT /Helvetica 8 Tf {self.PAGE_W / 2 - 25} {self.MARGIN_B - 24} Td ({safe_page}) Tj ET")

        # "Generated by EDISON" footer
        furniture.append(f"BT /Helvetica-Oblique 7 Tf {self.MARGIN_L} {self.MARGIN_B - 24} Td (Generated by EDISON AI) Tj ET")

        return furniture + page_content

    def render(self) -> bytes:
        """Render the document to PDF bytes."""
        # Finalize current page
        if self.current_page:
            self.pages.append(self.current_page)

        if not self.pages:
            self._new_page()
            self.add_paragraph("(Empty document)")
            self.pages.append(self.current_page)

        total_pages = len(self.pages)
        objects = []
        # Object 1: Catalog
        objects.append(b"1 0 obj << /Type /Catalog /Pages 2 0 R >> endobj")

        # Font objects — we need 4 fonts
        # Assign font object numbers after pages
        font_base = 3 + (2 * total_pages)
        font_names = ["Helvetica", "Helvetica-Bold", "Helvetica-Oblique", "Courier"]
        font_refs = ""
        for i, fname in enumerate(font_names):
            font_refs += f" /{fname} {font_base + i} 0 R"

        # Object 2: Pages
        kids = " ".join([f"{3 + i * 2} 0 R" for i in range(total_pages)])
        objects.append(f"2 0 obj << /Type /Pages /Kids [{kids}] /Count {total_pages} >> endobj".encode("ascii"))

        # Page + Content objects
        for idx, page_content in enumerate(self.pages):
            page_obj_num = 3 + idx * 2
            content_obj_num = 4 + idx * 2
            # Add page furniture (header/footer)
            full_content = self._add_page_furniture(page_content, idx + 1, total_pages)
            content_stream = "\n".join(full_content).encode("latin-1", errors="replace")
            page_obj = (
                f"{page_obj_num} 0 obj << /Type /Page /Parent 2 0 R "
                f"/MediaBox [0 0 {self.PAGE_W} {self.PAGE_H}] "
                f"/Contents {content_obj_num} 0 R "
                f"/Resources << /Font <<{font_refs}>> >> >> endobj"
            ).encode("ascii")
            content_obj = b"%d 0 obj << /Length %d >> stream\n%s\nendstream endobj" % (
                content_obj_num, len(content_stream), content_stream
            )
            objects.append(page_obj)
            objects.append(content_obj)

        # Font objects
        for i, fname in enumerate(font_names):
            obj_num = font_base + i
            objects.append(f"{obj_num} 0 obj << /Type /Font /Subtype /Type1 /BaseFont /{fname} >> endobj".encode("ascii"))

        # Build PDF
        xref_positions = []
        pdf = b"%PDF-1.4\n"
        for obj in objects:
            xref_positions.append(len(pdf))
            pdf += obj + b"\n"
        xref_start = len(pdf)
        total_objs = len(objects) + 1
        pdf += b"xref\n0 %d\n" % total_objs
        pdf += b"0000000000 65535 f \n"
        for pos in xref_positions:
            pdf += f"{pos:010d} 00000 n \n".encode("ascii")
        pdf += b"trailer << /Size %d /Root 1 0 R >>\nstartxref\n%d\n%%%%EOF" % (total_objs, xref_start)
        return pdf


def parse_markdown_to_pdf(markdown_text: str, title: str = "Document") -> bytes:
    """Convert markdown-formatted text to a professional PDF."""
    doc = PDFDocument(title=title)
    doc._new_page()

    lines = markdown_text.splitlines()
    i = 0
    in_code_block = False
    code_lines = []
    code_lang = ""
    in_table = False
    table_headers = []
    table_rows = []
    numbered_counter = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Code blocks
        if stripped.startswith("```"):
            if in_code_block:
                doc.add_code_block("\n".join(code_lines), code_lang)
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
                code_lang = stripped[3:].strip()
            i += 1
            continue

        if in_code_block:
            code_lines.append(line)
            i += 1
            continue

        # Tables
        if "|" in stripped and stripped.startswith("|"):
            cells = [c.strip() for c in stripped.split("|")[1:-1]]
            if all(re.match(r"^[-:]+$", c) for c in cells):
                # Separator row, skip
                i += 1
                continue
            if not in_table:
                in_table = True
                table_headers = cells
            else:
                table_rows.append(cells)
            i += 1
            continue
        elif in_table:
            doc.add_table(table_headers, table_rows)
            table_headers = []
            table_rows = []
            in_table = False

        # Empty line
        if not stripped:
            numbered_counter = 0
            doc.add_spacer(6)
            i += 1
            continue

        # Title (# heading)
        if stripped.startswith("# ") and not stripped.startswith("## "):
            doc.add_title(stripped[2:])
            i += 1
            continue

        # Heading 2
        if stripped.startswith("## "):
            doc.add_heading(stripped[3:], 1)
            i += 1
            continue

        # Heading 3
        if stripped.startswith("### "):
            doc.add_heading(stripped[4:], 2)
            i += 1
            continue

        # Heading 4+
        h_match = re.match(r"^(#{4,6})\s+(.+)$", stripped)
        if h_match:
            doc.add_heading(h_match.group(2), 2)
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^[-*_]{3,}$", stripped):
            doc.add_spacer(4)
            doc._draw_line(doc.MARGIN_L, doc.y, doc.PAGE_W - doc.MARGIN_R, doc.y, 0.5, (0.8, 0.8, 0.85))
            doc.y -= 8
            i += 1
            continue

        # Callout / blockquote
        if stripped.startswith("> "):
            quote_text = stripped[2:]
            style = "info"
            if "warning" in quote_text.lower() or "⚠" in quote_text:
                style = "warning"
            elif "note" in quote_text[:10].lower() or "ℹ" in quote_text:
                style = "info"
            elif "important" in quote_text.lower() or "❗" in quote_text:
                style = "error"
            doc.add_callout(quote_text, style)
            i += 1
            continue

        # Bullet points
        if re.match(r"^[-*+]\s+", stripped):
            text = re.sub(r"^[-*+]\s+", "", stripped)
            # Remove markdown bold/italic
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
            text = re.sub(r"\*(.+?)\*", r"\1", text)
            doc.add_bullet(text)
            i += 1
            continue

        # Sub-bullets
        if re.match(r"^\s{2,}[-*+]\s+", line):
            text = re.sub(r"^\s+[-*+]\s+", "", line)
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
            doc.add_bullet(text, indent=1)
            i += 1
            continue

        # Numbered items
        num_match = re.match(r"^(\d+)[.)]\s+(.+)$", stripped)
        if num_match:
            numbered_counter += 1
            text = num_match.group(2)
            text = re.sub(r"\*\*(.+?)\*\*", r"\1", text)
            doc.add_numbered_item(numbered_counter, text)
            i += 1
            continue

        # Regular paragraph — clean markdown formatting
        text = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
        text = re.sub(r"\*(.+?)\*", r"\1", text)
        text = re.sub(r"`(.+?)`", r"\1", text)
        text = re.sub(r"\[(.+?)\]\(.+?\)", r"\1", text)
        doc.add_paragraph(text)
        i += 1

    # Flush remaining table
    if in_table and table_headers:
        doc.add_table(table_headers, table_rows)

    # Flush remaining code
    if code_lines:
        doc.add_code_block("\n".join(code_lines), code_lang)

    return doc.render()


# ──────────────────────────────────────────────────────────────────
#  HTML Document Generation — Professional standalone HTML files
# ──────────────────────────────────────────────────────────────────

def generate_professional_html(content: str, title: str = "Document", style: str = "report") -> str:
    """Generate a professional HTML document with modern styling."""

    styles = {
        "report": {
            "primary": "#3b5bdb",
            "bg": "#ffffff",
            "text": "#212529",
            "font": "'Inter', 'Segoe UI', system-ui, sans-serif",
        },
        "presentation": {
            "primary": "#7950f2",
            "bg": "#1a1b2e",
            "text": "#e9ecef",
            "font": "'Inter', 'Segoe UI', system-ui, sans-serif",
        },
        "memo": {
            "primary": "#2b8a3e",
            "bg": "#ffffff",
            "text": "#212529",
            "font": "'Georgia', 'Times New Roman', serif",
        },
    }
    s = styles.get(style, styles["report"])
    date_str = datetime.now().strftime("%B %d, %Y")

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{html.escape(title)}</title>
    <link rel="preconnect" href="https://fonts.googleapis.com">
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700&display=swap" rel="stylesheet">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: {s['font']};
            background: {s['bg']};
            color: {s['text']};
            line-height: 1.7;
            font-size: 16px;
            -webkit-font-smoothing: antialiased;
        }}
        .container {{
            max-width: 800px;
            margin: 0 auto;
            padding: 60px 40px;
        }}
        .header {{
            border-bottom: 3px solid {s['primary']};
            padding-bottom: 24px;
            margin-bottom: 40px;
        }}
        .header h1 {{
            font-size: 2.2em;
            font-weight: 700;
            color: {s['primary']};
            margin-bottom: 8px;
            letter-spacing: -0.5px;
        }}
        .header .meta {{
            color: #868e96;
            font-size: 0.9em;
        }}
        h2 {{
            font-size: 1.5em;
            font-weight: 600;
            color: {s['primary']};
            margin: 32px 0 16px;
            padding-bottom: 8px;
            border-bottom: 1px solid #e9ecef;
        }}
        h3 {{
            font-size: 1.2em;
            font-weight: 600;
            margin: 24px 0 12px;
        }}
        p {{
            margin-bottom: 16px;
        }}
        ul, ol {{
            padding-left: 24px;
            margin-bottom: 16px;
        }}
        li {{
            margin-bottom: 8px;
        }}
        table {{
            width: 100%;
            border-collapse: collapse;
            margin: 24px 0;
            font-size: 0.95em;
        }}
        th {{
            background: {s['primary']};
            color: white;
            padding: 12px 16px;
            text-align: left;
            font-weight: 600;
        }}
        td {{
            padding: 10px 16px;
            border-bottom: 1px solid #e9ecef;
        }}
        tr:nth-child(even) {{
            background: #f8f9fa;
        }}
        blockquote {{
            border-left: 4px solid {s['primary']};
            padding: 16px 20px;
            margin: 20px 0;
            background: #f8f9fa;
            border-radius: 0 8px 8px 0;
            font-style: italic;
        }}
        code {{
            background: #f1f3f5;
            padding: 2px 6px;
            border-radius: 4px;
            font-family: 'Fira Code', 'Consolas', monospace;
            font-size: 0.9em;
        }}
        pre {{
            background: #212529;
            color: #f8f9fa;
            padding: 20px;
            border-radius: 8px;
            overflow-x: auto;
            margin: 20px 0;
            font-size: 0.9em;
            line-height: 1.5;
        }}
        pre code {{
            background: none;
            padding: 0;
            color: inherit;
        }}
        .callout {{
            padding: 16px 20px;
            border-radius: 8px;
            margin: 20px 0;
            border-left: 4px solid;
        }}
        .callout-info {{
            background: #e7f5ff;
            border-color: #339af0;
        }}
        .callout-warning {{
            background: #fff9db;
            border-color: #fcc419;
        }}
        .callout-success {{
            background: #ebfbee;
            border-color: #51cf66;
        }}
        hr {{
            border: none;
            border-top: 1px solid #dee2e6;
            margin: 32px 0;
        }}
        .footer {{
            margin-top: 60px;
            padding-top: 20px;
            border-top: 1px solid #dee2e6;
            color: #868e96;
            font-size: 0.85em;
            text-align: center;
        }}
        @media print {{
            body {{ font-size: 12pt; }}
            .container {{ padding: 0; max-width: none; }}
            pre {{ white-space: pre-wrap; }}
        }}
        @media (max-width: 600px) {{
            .container {{ padding: 24px 16px; }}
            h1 {{ font-size: 1.8em; }}
        }}
    </style>
</head>
<body>
    <div class="container">
        <div class="header">
            <h1>{html.escape(title)}</h1>
            <div class="meta">Generated on {date_str} · EDISON AI</div>
        </div>
        {content}
        <div class="footer">
            <p>Generated by EDISON AI · {date_str}</p>
        </div>
    </div>
</body>
</html>"""


# ──────────────────────────────────────────────────────────────────
#  Slideshow / Presentation Generation — HTML-based slide deck
# ──────────────────────────────────────────────────────────────────

def generate_slideshow_html(slides: List[Dict[str, Any]], title: str = "Presentation") -> str:
    """Generate a professional HTML slideshow presentation.
    
    Each slide dict can have:
        - title: str
        - content: str (HTML/markdown)
        - bullets: List[str]
        - layout: str (title, content, two-column, image)
        - notes: str (speaker notes)
    """
    date_str = datetime.now().strftime("%B %d, %Y")
    
    slides_html = []
    for i, slide in enumerate(slides):
        slide_title = html.escape(slide.get("title", ""))
        content = slide.get("content", "")
        bullets = slide.get("bullets", [])
        layout = slide.get("layout", "content")
        
        bullet_html = ""
        if bullets:
            bullet_html = "<ul class='slide-bullets'>" + "".join(
                f"<li>{html.escape(b)}</li>" for b in bullets
            ) + "</ul>"
        
        slide_class = f"slide slide-{layout}"
        if i == 0:
            slide_class += " active"
        
        slides_html.append(f"""
        <div class="{slide_class}" data-slide="{i}">
            <div class="slide-content">
                {"<h1>" + slide_title + "</h1>" if layout == "title" else "<h2>" + slide_title + "</h2>" if slide_title else ""}
                {f'<div class="slide-body">{content}</div>' if content else ""}
                {bullet_html}
            </div>
            <div class="slide-number">{i + 1} / {len(slides)}</div>
        </div>""")
    
    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{html.escape(title)}</title>
    <link href="https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap" rel="stylesheet">
    <style>
        * {{ margin: 0; padding: 0; box-sizing: border-box; }}
        body {{
            font-family: 'Inter', system-ui, sans-serif;
            background: #0f0f1a;
            color: #ffffff;
            overflow: hidden;
            height: 100vh;
            -webkit-font-smoothing: antialiased;
        }}
        .presentation {{
            width: 100vw;
            height: 100vh;
            position: relative;
        }}
        .slide {{
            position: absolute;
            top: 0; left: 0;
            width: 100%; height: 100%;
            display: flex;
            align-items: center;
            justify-content: center;
            padding: 60px 80px;
            opacity: 0;
            transform: translateX(100px);
            transition: all 0.5s cubic-bezier(0.4, 0, 0.2, 1);
            pointer-events: none;
            background: linear-gradient(135deg, #0f0f1a 0%, #1a1a2e 50%, #16213e 100%);
        }}
        .slide.active {{
            opacity: 1;
            transform: translateX(0);
            pointer-events: all;
        }}
        .slide.prev {{
            transform: translateX(-100px);
        }}
        .slide-content {{
            max-width: 1000px;
            width: 100%;
        }}
        .slide-title .slide-content {{
            text-align: center;
        }}
        .slide h1 {{
            font-size: 3.5em;
            font-weight: 800;
            background: linear-gradient(135deg, #667eea, #764ba2);
            -webkit-background-clip: text;
            -webkit-text-fill-color: transparent;
            margin-bottom: 24px;
            letter-spacing: -1px;
            line-height: 1.1;
        }}
        .slide h2 {{
            font-size: 2em;
            font-weight: 700;
            color: #667eea;
            margin-bottom: 32px;
            padding-bottom: 16px;
            border-bottom: 2px solid rgba(102, 126, 234, 0.3);
        }}
        .slide-body {{
            font-size: 1.3em;
            line-height: 1.8;
            color: #c4c4d4;
        }}
        .slide-bullets {{
            list-style: none;
            padding: 0;
        }}
        .slide-bullets li {{
            font-size: 1.3em;
            line-height: 1.6;
            color: #c4c4d4;
            padding: 12px 0 12px 40px;
            position: relative;
            border-bottom: 1px solid rgba(255,255,255,0.05);
        }}
        .slide-bullets li::before {{
            content: '▸';
            position: absolute;
            left: 10px;
            color: #667eea;
            font-size: 1.2em;
        }}
        .slide-number {{
            position: absolute;
            bottom: 24px;
            right: 40px;
            font-size: 0.9em;
            color: rgba(255,255,255,0.3);
            font-weight: 500;
        }}
        .controls {{
            position: fixed;
            bottom: 24px;
            left: 50%;
            transform: translateX(-50%);
            display: flex;
            gap: 12px;
            z-index: 100;
        }}
        .controls button {{
            background: rgba(255,255,255,0.1);
            border: 1px solid rgba(255,255,255,0.2);
            color: white;
            padding: 10px 24px;
            border-radius: 8px;
            cursor: pointer;
            font-family: inherit;
            font-size: 14px;
            transition: all 0.2s;
        }}
        .controls button:hover {{
            background: rgba(102, 126, 234, 0.3);
            border-color: #667eea;
        }}
        .progress {{
            position: fixed;
            top: 0;
            left: 0;
            height: 3px;
            background: linear-gradient(90deg, #667eea, #764ba2);
            transition: width 0.3s ease;
            z-index: 200;
        }}
        @media (max-width: 768px) {{
            .slide {{ padding: 40px 24px; }}
            .slide h1 {{ font-size: 2em; }}
            .slide h2 {{ font-size: 1.4em; }}
            .slide-bullets li {{ font-size: 1em; }}
        }}
        @media print {{
            .slide {{
                position: relative;
                opacity: 1;
                transform: none;
                pointer-events: all;
                page-break-after: always;
                height: 100vh;
            }}
            .controls {{ display: none; }}
            .progress {{ display: none; }}
        }}
    </style>
</head>
<body>
    <div class="progress" id="progress"></div>
    <div class="presentation" id="presentation">
        {''.join(slides_html)}
    </div>
    <div class="controls">
        <button onclick="prevSlide()">← Previous</button>
        <button onclick="nextSlide()">Next →</button>
    </div>
    <script>
        let current = 0;
        const slides = document.querySelectorAll('.slide');
        const progress = document.getElementById('progress');
        
        function updateProgress() {{
            const pct = ((current + 1) / slides.length) * 100;
            progress.style.width = pct + '%';
        }}
        
        function showSlide(n) {{
            slides.forEach((s, i) => {{
                s.classList.remove('active', 'prev');
                if (i === n) s.classList.add('active');
                else if (i < n) s.classList.add('prev');
            }});
            current = n;
            updateProgress();
        }}
        
        function nextSlide() {{ if (current < slides.length - 1) showSlide(current + 1); }}
        function prevSlide() {{ if (current > 0) showSlide(current - 1); }}
        
        document.addEventListener('keydown', (e) => {{
            if (e.key === 'ArrowRight' || e.key === ' ') nextSlide();
            if (e.key === 'ArrowLeft') prevSlide();
        }});
        
        // Touch support
        let startX = 0;
        document.addEventListener('touchstart', e => startX = e.touches[0].clientX);
        document.addEventListener('touchend', e => {{
            const diff = e.changedTouches[0].clientX - startX;
            if (Math.abs(diff) > 50) {{
                if (diff < 0) nextSlide();
                else prevSlide();
            }}
        }});
        
        updateProgress();
    </script>
</body>
</html>"""


# ──────────────────────────────────────────────────────────────────
#  CSV / Spreadsheet Generation
# ──────────────────────────────────────────────────────────────────

def generate_csv(headers: List[str], rows: List[List[str]]) -> str:
    """Generate a properly formatted CSV string."""
    import csv
    import io
    output = io.StringIO()
    writer = csv.writer(output, quoting=csv.QUOTE_MINIMAL)
    writer.writerow(headers)
    for row in rows:
        writer.writerow(row)
    return output.getvalue()


# ──────────────────────────────────────────────────────────────────
#  Rich Text / Markdown Document
# ──────────────────────────────────────────────────────────────────

def generate_rich_markdown(title: str, sections: List[Dict[str, Any]]) -> str:
    """Generate a well-structured Markdown document."""
    date_str = datetime.now().strftime("%B %d, %Y")
    parts = [
        f"# {title}\n",
        f"*Generated on {date_str} by EDISON AI*\n",
        "---\n",
    ]
    
    # Table of contents
    parts.append("## Table of Contents\n")
    for i, sec in enumerate(sections, 1):
        sec_title = sec.get("title", f"Section {i}")
        anchor = sec_title.lower().replace(" ", "-").replace(":", "")
        parts.append(f"{i}. [{sec_title}](#{anchor})")
    parts.append("\n---\n")
    
    for sec in sections:
        sec_title = sec.get("title", "Section")
        content = sec.get("content", "")
        parts.append(f"## {sec_title}\n")
        parts.append(f"{content}\n")
    
    parts.append("\n---\n")
    parts.append(f"*End of document · {date_str}*\n")
    
    return "\n".join(parts)


# ──────────────────────────────────────────────────────────────────
#  File Generation Dispatch — Unified entry point
# ──────────────────────────────────────────────────────────────────

# ──────────────────────────────────────────────────────────────────
#  PPTX (PowerPoint) Generation — Real .pptx files via python-pptx
# ──────────────────────────────────────────────────────────────────

def generate_pptx(slides: List[Dict[str, Any]], title: str = "Presentation") -> bytes:
    """Generate a real .pptx PowerPoint file from slide data.
    
    Each slide dict can have:
        - title: str
        - content: str (plain text paragraph)
        - bullets: List[str]
        - layout: str (title, content, two-column)
        - notes: str (speaker notes)
    """
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        logger.warning("python-pptx not installed, falling back to HTML slideshow")
        html_content = generate_slideshow_html(slides, title)
        return html_content.encode("utf-8"), ".html"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color scheme
    primary_rgb = RGBColor(0x66, 0x7E, 0xEA)
    dark_bg = RGBColor(0x0F, 0x0F, 0x1A)
    light_text = RGBColor(0xE0, 0xE0, 0xF0)
    muted_text = RGBColor(0xC4, 0xC4, 0xD4)

    for i, slide_data in enumerate(slides):
        slide_title = slide_data.get("title", "")
        content = slide_data.get("content", "")
        bullets = slide_data.get("bullets", [])
        layout = slide_data.get("layout", "content")
        notes = slide_data.get("notes", "")

        # Use blank layout and build manually for consistent styling
        blank_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(blank_layout)

        # Dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = dark_bg

        if layout == "title":
            # Title slide - centered title and subtitle
            title_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(2.0), Inches(10.333), Inches(2.0)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(44)
            p.font.bold = True
            p.font.color.rgb = primary_rgb
            p.alignment = PP_ALIGN.CENTER

            if content:
                subtitle_box = slide.shapes.add_textbox(
                    Inches(2.0), Inches(4.2), Inches(9.333), Inches(1.5)
                )
                tf2 = subtitle_box.text_frame
                tf2.word_wrap = True
                p2 = tf2.paragraphs[0]
                p2.text = content
                p2.font.size = Pt(20)
                p2.font.color.rgb = muted_text
                p2.alignment = PP_ALIGN.CENTER
        else:
            # Content slide
            # Title at top
            title_box = slide.shapes.add_textbox(
                Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.0)
            )
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = primary_rgb

            # Accent line under title
            line = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.8), Inches(1.55), Inches(3), Pt(3)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = primary_rgb
            line.line.fill.background()

            # Content area
            content_top = Inches(1.8)
            content_box = slide.shapes.add_textbox(
                Inches(0.8), content_top, Inches(11.733), Inches(5.0)
            )
            tf_content = content_box.text_frame
            tf_content.word_wrap = True

            if bullets:
                for j, bullet in enumerate(bullets):
                    if j == 0:
                        p = tf_content.paragraphs[0]
                    else:
                        p = tf_content.add_paragraph()
                    p.text = f"  \u25b8  {bullet}"
                    p.font.size = Pt(20)
                    p.font.color.rgb = light_text
                    p.space_after = Pt(12)
            elif content:
                p = tf_content.paragraphs[0]
                p.text = content
                p.font.size = Pt(18)
                p.font.color.rgb = light_text
                p.line_spacing = Pt(28)

        # Slide number
        num_box = slide.shapes.add_textbox(
            Inches(12.0), Inches(6.8), Inches(1.2), Inches(0.5)
        )
        num_tf = num_box.text_frame
        num_p = num_tf.paragraphs[0]
        num_p.text = f"{i + 1} / {len(slides)}"
        num_p.font.size = Pt(10)
        num_p.font.color.rgb = RGBColor(0x80, 0x80, 0x90)
        num_p.alignment = PP_ALIGN.RIGHT

        # Speaker notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

    import io
    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()


# ──────────────────────────────────────────────────────────────────
#  DOCX (Word) Generation — Real .docx files via python-docx
# ──────────────────────────────────────────────────────────────────

def _clean_markdown_text(text: str) -> str:
    """Strip markdown formatting from text to produce clean prose."""
    if not text:
        return ""
    # Remove markdown bold/italic
    text = re.sub(r'\*\*\*(.+?)\*\*\*', r'\1', text)
    text = re.sub(r'\*\*(.+?)\*\*', r'\1', text)
    text = re.sub(r'\*(.+?)\*', r'\1', text)
    text = re.sub(r'__(.+?)__', r'\1', text)
    text = re.sub(r'_(.+?)_', r'\1', text)
    # Remove inline code backticks
    text = re.sub(r'`([^`]+)`', r'\1', text)
    # Remove link syntax [text](url) → text
    text = re.sub(r'\[(.+?)\]\(.+?\)', r'\1', text)
    # Remove image syntax ![alt](url)
    text = re.sub(r'!\[.*?\]\(.+?\)', '', text)
    return text.strip()


def generate_docx(content: str, title: str = "Document") -> bytes:
    """Generate a real .docx Word document from markdown-formatted content.
    
    Parses markdown and converts it to properly formatted Word paragraphs
    with headings, bullet points, numbered lists, and clean body text.
    """
    try:
        from docx import Document
        from docx.shared import Pt, Inches, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH
    except ImportError:
        logger.warning("python-docx not installed, falling back to plain text")
        return content.encode("utf-8"), ".txt"

    doc = Document()

    # Set default font style
    style = doc.styles['Normal']
    font = style.font
    font.name = 'Calibri'
    font.size = Pt(11)
    font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

    # Title
    title_para = doc.add_heading(title, level=0)
    title_para.alignment = WD_ALIGN_PARAGRAPH.LEFT

    # Date subtitle
    date_str = datetime.now().strftime("%B %d, %Y")
    subtitle = doc.add_paragraph(f"Generated on {date_str}")
    subtitle.style = doc.styles['Subtitle']

    doc.add_paragraph()  # Spacer

    # Parse markdown content into the document
    lines = content.splitlines()
    i = 0
    in_code_block = False
    code_lines = []
    numbered_counter = 0

    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        # Code blocks
        if stripped.startswith("```"):
            if in_code_block:
                # End code block - add as a styled paragraph
                code_text = "\n".join(code_lines)
                code_para = doc.add_paragraph()
                code_run = code_para.add_run(code_text)
                code_run.font.name = 'Consolas'
                code_run.font.size = Pt(9)
                code_run.font.color.rgb = RGBColor(0x20, 0x20, 0x25)
                code_lines = []
                in_code_block = False
            else:
                in_code_block = True
            i += 1
            continue

        if in_code_block:
            code_lines.append(line)
            i += 1
            continue

        # Empty line
        if not stripped:
            numbered_counter = 0
            i += 1
            continue

        # Headings
        if stripped.startswith("# ") and not stripped.startswith("## "):
            doc.add_heading(_clean_markdown_text(stripped[2:]), level=1)
            i += 1
            continue
        if stripped.startswith("## "):
            doc.add_heading(_clean_markdown_text(stripped[3:]), level=2)
            i += 1
            continue
        if stripped.startswith("### "):
            doc.add_heading(_clean_markdown_text(stripped[4:]), level=3)
            i += 1
            continue
        h_match = re.match(r"^(#{4,6})\s+(.+)$", stripped)
        if h_match:
            doc.add_heading(_clean_markdown_text(h_match.group(2)), level=4)
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^[-*_]{3,}$", stripped):
            doc.add_paragraph("─" * 50)
            i += 1
            continue

        # Blockquote
        if stripped.startswith("> "):
            quote_text = _clean_markdown_text(stripped[2:])
            quote_para = doc.add_paragraph(quote_text)
            quote_para.style = doc.styles['Intense Quote'] if 'Intense Quote' in doc.styles else doc.styles['Quote']
            i += 1
            continue

        # Bullet points
        if re.match(r'^[-*+]\s+', stripped):
            text = re.sub(r'^[-*+]\s+', '', stripped)
            doc.add_paragraph(_clean_markdown_text(text), style='List Bullet')
            i += 1
            continue

        # Sub-bullets
        if re.match(r'^\s{2,}[-*+]\s+', line):
            text = re.sub(r'^\s+[-*+]\s+', '', line)
            doc.add_paragraph(_clean_markdown_text(text), style='List Bullet 2')
            i += 1
            continue

        # Numbered items
        num_match = re.match(r'^(\d+)[.)]\s+(.+)$', stripped)
        if num_match:
            text = num_match.group(2)
            doc.add_paragraph(_clean_markdown_text(text), style='List Number')
            i += 1
            continue

        # Regular paragraph - clean all markdown formatting
        clean = _clean_markdown_text(stripped)
        if clean:
            doc.add_paragraph(clean)
        i += 1

    # Flush remaining code block
    if code_lines:
        code_text = "\n".join(code_lines)
        code_para = doc.add_paragraph()
        code_run = code_para.add_run(code_text)
        code_run.font.name = 'Consolas'
        code_run.font.size = Pt(9)

    # Footer
    doc.add_paragraph()
    footer = doc.add_paragraph(f"Generated by EDISON AI — {date_str}")
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer.runs[0]
    footer_run.font.size = Pt(8)
    footer_run.font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    buffer = io.BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


FILE_GENERATION_PROMPT = """When the user asks you to create a downloadable file, you MUST output a FILES block. Follow these rules carefully:

## IMPORTANT: What NOT to create as files
- **NEVER create audio files** (.mp3, .wav, .ogg, .flac, .aac, .m4a) — use the generate_music tool instead
- **NEVER create video files** (.mp4, .avi, .mov, .webm, .mkv) — use the generate_video tool instead
- **NEVER create image files** (.png, .jpg, .jpeg, .gif, .webp) — use the generate_image tool instead
- If the user asks for music, songs, beats, audio, or sound, call generate_music — do NOT create a file
- If the user asks for video, animation, or clips, call generate_video — do NOT create a file

## File Types You Can Create
- **PDF documents**: Reports, essays, papers, guides, manuals, resumes, letters
- **Word documents (.docx)**: Essays, research papers, reports, formal documents
- **Presentations/Slideshows (.pptx)**: PowerPoint slide decks with multiple slides
- **HTML documents**: Rich formatted interactive pages
- **Spreadsheets (CSV)**: Data tables, lists, organized data
- **Text documents (TXT/MD)**: Notes, outlines, plain text files
- **Code files**: Scripts, source code in any language
- **JSON data**: Structured data, configurations
- **ZIP archives**: Multiple files bundled together

## Output Format

For essays, research papers, reports, and written documents use .docx:
```files
[{"filename": "document.docx", "content": "# Title\\n\\nFull detailed content with markdown formatting. Use ## for sections, - for bullets, 1. for numbered lists. Write complete paragraphs of clean prose."}]
```

For PDF documents (resumes, formal letters, technical specs):
```files
[{"filename": "report.pdf", "content": "# Title\\n\\nFull detailed content with markdown formatting..."}]
```

For presentations/slideshows use .pptx (NOT .html):
```files
[{"filename": "presentation.pptx", "content": "", "type": "slideshow", "slides": [
  {"title": "Presentation Title", "content": "Subtitle or description", "layout": "title"},
  {"title": "Key Points", "bullets": ["First important point", "Second point with detail", "Third compelling point"], "layout": "content"},
  {"title": "Details", "content": "Plain text content for this slide.", "layout": "content"},
  {"title": "Thank You", "content": "Questions and contact info", "layout": "title"}
]}]
```

For other files (TXT, MD, CSV, JSON, code files):
```files
[{"filename": "data.csv", "content": "col1,col2\\nval1,val2"}]
```

## CRITICAL Rules
1. **Be comprehensive**: Write FULL, DETAILED content. Never use placeholders like "[Add more here]"
2. **Use proper structure**: Include headings, paragraphs, bullet points, numbered lists
3. **Content is markdown**: It gets converted to the target format with proper styling. Use # headings, - bullets, **bold**, etc.
4. **Presentations**: Use .pptx extension. Create 5-10 slides minimum, each with meaningful content
5. **Word documents**: Use .docx extension for essays, research papers, reports
6. **Do NOT repeat content**: Write each section once. Do NOT loop or duplicate paragraphs
7. **Keep the ```files block compact**: Put all file content in the JSON. Add only a brief 1-sentence summary outside the block

Keep your summary outside the files block to just one brief sentence like "Here's your document."
Do NOT repeat the file content outside the block."""


def render_file_entry(entry: dict) -> Tuple[str, bytes]:
    """
    Render a file entry from LLM output into final filename + bytes.
    Supports enhanced types: slideshow/pptx, docx, html_document, and auto-detects
    markdown content for PDF conversion.
    
    Returns (filename, data_bytes)
    """
    filename = entry.get("filename", "output.txt")
    content = entry.get("content", "")
    file_type = entry.get("type", "").lower()
    ext = os.path.splitext(filename)[1].lower()

    # ─── Slideshow / Presentation ───
    if file_type == "slideshow" or ext in (".pptx",):
        slides = entry.get("slides", [])
        if not slides:
            # Try to parse slides from content
            slides = [{"title": "Presentation", "content": str(content), "layout": "title"}]
        title = entry.get("title", os.path.splitext(filename)[0])
        try:
            pptx_data = generate_pptx(slides, title)
            if isinstance(pptx_data, tuple):
                # Fallback returned (data, ext)
                return os.path.splitext(filename)[0] + pptx_data[1], pptx_data[0]
            # Force .pptx extension
            if not filename.endswith(".pptx"):
                filename = os.path.splitext(filename)[0] + ".pptx"
            return filename, pptx_data
        except Exception as e:
            logger.warning(f"PPTX generation failed, falling back to HTML: {e}")
            html_content = generate_slideshow_html(slides, title)
            if not filename.endswith(".html"):
                filename = os.path.splitext(filename)[0] + ".html"
            return filename, html_content.encode("utf-8")

    # ─── Word Document (.docx) ───
    if ext == ".docx" or ext == ".doc":
        title = entry.get("title", os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title())
        try:
            docx_data = generate_docx(str(content), title)
            if isinstance(docx_data, tuple):
                return os.path.splitext(filename)[0] + docx_data[1], docx_data[0]
            if not filename.endswith(".docx"):
                filename = os.path.splitext(filename)[0] + ".docx"
            return filename, docx_data
        except Exception as e:
            logger.warning(f"DOCX generation failed, falling back to PDF: {e}")
            data = parse_markdown_to_pdf(str(content), title)
            filename = os.path.splitext(filename)[0] + ".pdf"
            return filename, data

    # ─── Rich HTML Document ───
    if file_type == "html_document":
        title = entry.get("title", os.path.splitext(filename)[0])
        html_content = generate_professional_html(str(content), title)
        if not filename.endswith(".html"):
            filename = os.path.splitext(filename)[0] + ".html"
        return filename, html_content.encode("utf-8")

    # ─── PDF ───
    if ext == ".pdf":
        title = entry.get("title", os.path.splitext(filename)[0].replace("_", " ").replace("-", " ").title())
        data = parse_markdown_to_pdf(str(content), title)
        return filename, data

    # ─── CSV ───
    if ext == ".csv" and isinstance(content, str):
        return filename, content.encode("utf-8")

    # ─── JSON ───
    if ext == ".json":
        if isinstance(content, (dict, list)):
            return filename, json.dumps(content, indent=2).encode("utf-8")
        return filename, str(content).encode("utf-8")

    # ─── Default text ───
    if isinstance(content, (dict, list)):
        return filename, json.dumps(content, indent=2).encode("utf-8")
    return filename, str(content).encode("utf-8")
